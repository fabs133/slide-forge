"""Tests for slideforge.renderer."""
from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation as PptxPresentation

from slideforge.models import Presentation, Slide
from slideforge.renderer import render_pptx
from slideforge.template_loader import TemplateLayoutError

TEMPLATE_DIR = Path(__file__).parent.parent / "slideforge" / "templates" / "default"
pytestmark = pytest.mark.skipif(not TEMPLATE_DIR.exists(), reason="Template not generated")


def test_render_creates_file(tmp_path, sample_presentation):
    """Tests the rendering of a presentation to a file."""
    out = tmp_path / "output.pptx"
    render_pptx(sample_presentation, out)
    assert out.exists()
    assert out.stat().st_size > 0


def test_render_correct_slide_count(tmp_path, sample_presentation):
    """Tests the correct number of slides are created."""
    out = tmp_path / "output.pptx"
    render_pptx(sample_presentation, out)
    prs = PptxPresentation(str(out))
    assert len(prs.slides) == 3


def test_render_title_text(tmp_path, sample_presentation):
    """Tests rendering of title text in a presentation."""
    out = tmp_path / "output.pptx"
    render_pptx(sample_presentation, out)
    prs = PptxPresentation(str(out))
    first_slide = prs.slides[0]
    texts = [sh.text for sh in first_slide.shapes if sh.has_text_frame]
    assert any("Welcome" in t for t in texts)


def test_render_body_bullets(tmp_path):
    """Render body bullets for a presentation."""
    pres = Presentation(id="bullets", name="Bullets", slides=[
        Slide(id="b1", layout="SP_Content", title="Pts", body="- Point A\n- Point B\n- Point C"),
    ])
    out = tmp_path / "output.pptx"
    render_pptx(pres, out)
    prs = PptxPresentation(str(out))
    slide = prs.slides[0]
    body_texts = []
    for sh in slide.shapes:
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                if p.text:
                    body_texts.append(p.text)
    assert "Point A" in body_texts
    assert "Point B" in body_texts


def test_render_section_break_no_body(tmp_path):
    """Test rendering a section break slide without body text."""
    pres = Presentation(id="sb", name="SB", slides=[
        Slide(id="sb1", layout="SP_SectionBreak", title="Chapter 2", body="This should not appear"),
    ])
    out = tmp_path / "output.pptx"
    render_pptx(pres, out)
    prs = PptxPresentation(str(out))
    slide = prs.slides[0]
    all_text = " ".join(sh.text for sh in slide.shapes if sh.has_text_frame)
    assert "Chapter 2" in all_text
    assert "should not appear" not in all_text


def test_render_speaker_notes(tmp_path):
    """Test rendering of speaker notes."""
    pres = Presentation(id="notes", name="Notes", slides=[
        Slide(id="n1", layout="SP_Content", title="T", body="B", notes="Remember to explain X"),
    ])
    out = tmp_path / "output.pptx"
    render_pptx(pres, out)
    prs = PptxPresentation(str(out))
    notes_text = prs.slides[0].notes_slide.notes_text_frame.text
    assert "Remember to explain X" in notes_text


def test_render_empty_presentation(tmp_path):
    """Test rendering an empty presentation."""
    pres = Presentation(id="empty", name="Empty")
    out = tmp_path / "output.pptx"
    render_pptx(pres, out)
    prs = PptxPresentation(str(out))
    assert len(prs.slides) == 0


def test_render_unknown_layout_raises(tmp_path):
    """Test rendering a presentation with an unknown layout raises an exception."""
    pres = Presentation(id="bad", name="Bad", slides=[
        Slide(id="x", layout="SP_Nonexistent", title="T"),
    ])
    out = tmp_path / "output.pptx"
    with pytest.raises(TemplateLayoutError):
        render_pptx(pres, out)


def test_render_all_layouts(tmp_path):
    """Render one slide of each layout type to verify none crash."""
    layouts = ["SP_Title", "SP_Content", "SP_Intro", "SP_Closing", "SP_Sources", "SP_SectionBreak", "SP_Code"]
    slides = [Slide(id=f"s{i}", layout=lay, title=f"Slide {i}", body=f"Body {i}") for i, lay in enumerate(layouts)]
    pres = Presentation(id="all", name="All Layouts", slides=slides)
    out = tmp_path / "output.pptx"
    render_pptx(pres, out)
    prs = PptxPresentation(str(out))
    assert len(prs.slides) == 7
