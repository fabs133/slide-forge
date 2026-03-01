"""
generate_template.py — Build the slide-forge PPTX template from code.

Produces:
  slideforge/templates/default/slides.pptx
  slideforge/templates/default/config.json

Run from the repo root:
  python -m slideforge.tools.generate_template

No PowerPoint, LibreOffice, or any GUI tool needed.
"""

from __future__ import annotations

import json
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.util import Cm, Pt

# ── Colour palette (clean preset) ─────────────────────────────────────────────


def _rgb(hex_str: str) -> RGBColor:
    """Converts a hexadecimal color string to an RGBColor object.

    :param hex_str: Hexadecimal color string starting with '#' and followed by 6 digits.
    :type hex_str: str
    :return: RGBColor object representing the given color.
    :rtype: RGBColor
    :raises ValueError: If the input string is not a valid hexadecimal color.
    """
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


NAVY = _rgb("#1E2761")
LIGHT = _rgb("#F5F7FA")
CORAL = _rgb("#F96167")
BLUE2 = _rgb("#CADCFC")
DARK = _rgb("#212121")
WHITE = _rgb("#FFFFFF")

# Slide canvas
W = Cm(33.87)
H = Cm(19.05)


# ── Layout configuration ───────────────────────────────────────────────────────

LAYOUTS = [
    {
        "layout_idx": 0,  # "Title Slide"  → ctrTitle(0) + subTitle(1)
        "name": "SP_Title",
        "bg": NAVY,
        "title": (Cm(1.5), Cm(5.5), Cm(26), Cm(4), Pt(44), True, WHITE),
        "body": (Cm(1.5), Cm(10.5), Cm(22), Cm(2), Pt(22), False, BLUE2),
        "extra_phs": [],
        "deco": [
            (Cm(0), Cm(0), Cm(0.4), H, CORAL),  # left coral bar
        ],
    },
    {
        "layout_idx": 1,  # "Title and Content"  → title(0) + body(1)
        "name": "SP_Content",
        "bg": LIGHT,
        "title": (Cm(2), Cm(1.0), Cm(29), Cm(1.5), Pt(32), True, NAVY),
        "body": (Cm(2), Cm(3.2), Cm(20), Cm(13), Pt(18), False, DARK),
        "extra_phs": [],
        "deco": [
            (Cm(0), Cm(0), W, Cm(0.35), NAVY),  # top bar
            (Cm(2), Cm(2.6), Cm(26), Cm(0.15), CORAL),  # heading underline
        ],
    },
    {
        "layout_idx": 2,  # "Section Header"  → title(0) + body(1)
        "name": "SP_Intro",
        "bg": LIGHT,
        "title": (Cm(2.5), Cm(1.0), Cm(28), Cm(1.5), Pt(30), True, NAVY),
        "body": (Cm(3.5), Cm(3.0), Cm(26), Cm(13.5), Pt(18), False, DARK),
        "extra_phs": [],
        "deco": [
            (Cm(0), Cm(0), W, Cm(0.35), NAVY),  # top bar
            (Cm(0.8), Cm(2.5), Cm(0.5), Cm(13), BLUE2),  # left accent
        ],
    },
    {
        "layout_idx": 3,  # "Two Content"  → title(0) + body(1) + extra body(2)
        "name": "SP_Closing",
        "bg": LIGHT,
        "title": (Cm(2), Cm(1.0), Cm(28), Cm(1.6), Pt(32), True, NAVY),
        "body": (Cm(2), Cm(3.5), Cm(28), Cm(12), Pt(18), False, DARK),
        "extra_phs": [2],  # remove the second content area
        "deco": [
            (Cm(0), Cm(0), W, Cm(0.35), NAVY),  # top bar
            (Cm(0), Cm(18.55), W, Cm(0.5), NAVY),  # bottom bar
            (Cm(2), Cm(2.7), Cm(6), Cm(0.2), CORAL),  # accent line
        ],
    },
    {
        "layout_idx": 4,  # "Comparison"  → title(0) + body(1) + extras(2,3,4)
        "name": "SP_Sources",
        "bg": NAVY,
        "title": (Cm(2), Cm(0.8), Cm(28), Cm(1.5), Pt(28), True, WHITE),
        "body": (Cm(2), Cm(2.8), Cm(29), Cm(14), Pt(13), False, BLUE2),
        "extra_phs": [2, 3, 4],  # remove all the comparison columns
        "deco": [],
    },
    {
        "layout_idx": 5,  # "Title Only"  → title(0), no body
        "name": "SP_SectionBreak",
        "bg": NAVY,
        "title": (Cm(1.5), Cm(6.5), Cm(15), Cm(5), Pt(40), True, WHITE),
        "body": None,  # intentionally no body placeholder
        "extra_phs": [],
        "deco": [
            (Cm(17.87), Cm(0), Cm(16), H, CORAL),  # right coral half
        ],
    },
    {
        "layout_idx": 6,  # "Blank" → repurposed for code slides
        "name": "SP_Code",
        "bg": LIGHT,
        "title": (Cm(2), Cm(1.0), Cm(29), Cm(1.5), Pt(28), True, NAVY),
        "body": (Cm(2), Cm(3.0), Cm(29), Cm(14), Pt(14), False, DARK),
        "extra_phs": [],
        "deco": [
            (Cm(0), Cm(0), W, Cm(0.35), NAVY),  # top bar
        ],
    },
]

# ── Helpers ────────────────────────────────────────────────────────────────────


def _set_ph_position(ph_elem, x, y, w, h) -> None:
    """Update the xfrm position/size of a placeholder element."""
    spPr = ph_elem.find(qn("p:spPr"))
    if spPr is None:
        spPr = etree.SubElement(ph_elem, qn("p:spPr"))

    xfrm = spPr.find(qn("a:xfrm"))
    if xfrm is None:
        xfrm = etree.SubElement(spPr, qn("a:xfrm"))

    for tag in ("a:off", "a:ext"):
        old = xfrm.find(qn(tag))
        if old is not None:
            xfrm.remove(old)

    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(int(x)))
    off.set("y", str(int(y)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(int(w)))
    ext.set("cy", str(int(h)))


def _set_ph_font(ph_elem, sz: Pt, bold: bool, color: RGBColor) -> None:
    """Set default run properties on a placeholder's lstStyle level 1."""
    txBody = ph_elem.find(qn("p:txBody"))
    if txBody is None:
        return
    lstStyle = txBody.find(qn("a:lstStyle"))
    if lstStyle is None:
        lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))

    for child in list(lstStyle):
        lstStyle.remove(child)

    lvl1 = etree.SubElement(lstStyle, qn("a:lvl1pPr"))
    defRPr = etree.SubElement(lvl1, qn("a:defRPr"))
    defRPr.set("sz", str(int(sz.pt * 100)))
    defRPr.set("b", "1" if bold else "0")
    solidFill = etree.SubElement(defRPr, qn("a:solidFill"))
    srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
    srgbClr.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")


def _remove_placeholder(spTree, idx: int) -> None:
    """Remove a placeholder by its idx value from spTree."""
    for sp in list(spTree.findall(qn("p:sp"))):
        nvPr = sp.find(f".//{qn('p:ph')}")
        if nvPr is not None and nvPr.get("idx") == str(idx):
            spTree.remove(sp)
            return


def _remove_date_footer_slidenum(spTree) -> None:
    """Remove the standard date/footer/slide-number placeholders (idx 10, 11, 12)."""
    for idx in (10, 11, 12):
        _remove_placeholder(spTree, idx)


def _add_rect_shape(spTree, elem_id: int, x, y, w, h, color: RGBColor) -> None:
    """Append a solid rectangle shape (decorative) to an spTree."""
    sp = etree.SubElement(spTree, qn("p:sp"))

    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qn("p:cNvPr"))
    cNvPr.set("id", str(elem_id))
    cNvPr.set("name", f"deco_{elem_id}")
    etree.SubElement(nvSpPr, qn("p:cNvSpPr"))
    etree.SubElement(nvSpPr, qn("p:nvPr"))

    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm, qn("a:off"))
    off.set("x", str(int(x)))
    off.set("y", str(int(y)))
    ext = etree.SubElement(xfrm, qn("a:ext"))
    ext.set("cx", str(int(w)))
    ext.set("cy", str(int(h)))
    prstGeom = etree.SubElement(spPr, qn("a:prstGeom"))
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, qn("a:avLst"))

    fill = etree.SubElement(spPr, qn("a:solidFill"))
    srgb = etree.SubElement(fill, qn("a:srgbClr"))
    srgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    ln = etree.SubElement(spPr, qn("a:ln"))
    etree.SubElement(ln, qn("a:noFill"))

    txBody = etree.SubElement(sp, qn("p:txBody"))
    etree.SubElement(txBody, qn("a:bodyPr"))
    etree.SubElement(txBody, qn("a:lstStyle"))
    etree.SubElement(txBody, qn("a:p"))


def _add_placeholder(spTree, idx: int, x, y, w, h) -> None:
    """Add a new placeholder shape to a layout that doesn't have one (e.g. Blank)."""
    sp = etree.SubElement(spTree, qn("p:sp"))

    nvSpPr = etree.SubElement(sp, qn("p:nvSpPr"))
    cNvPr = etree.SubElement(nvSpPr, qn("p:cNvPr"))
    cNvPr.set("id", str(900 + idx))
    cNvPr.set("name", f"Placeholder {idx}")
    cNvSpPr = etree.SubElement(nvSpPr, qn("p:cNvSpPr"))
    spLocks = etree.SubElement(cNvSpPr, qn("a:spLocks"))
    spLocks.set("noGrp", "1")
    nvPr = etree.SubElement(nvSpPr, qn("p:nvPr"))
    ph = etree.SubElement(nvPr, qn("p:ph"))
    ph.set("idx", str(idx))
    if idx == 0:
        ph.set("type", "title")

    spPr = etree.SubElement(sp, qn("p:spPr"))
    xfrm_el = etree.SubElement(spPr, qn("a:xfrm"))
    off = etree.SubElement(xfrm_el, qn("a:off"))
    off.set("x", str(int(x)))
    off.set("y", str(int(y)))
    ext = etree.SubElement(xfrm_el, qn("a:ext"))
    ext.set("cx", str(int(w)))
    ext.set("cy", str(int(h)))

    txBody = etree.SubElement(sp, qn("p:txBody"))
    etree.SubElement(txBody, qn("a:bodyPr"))
    lstStyle = etree.SubElement(txBody, qn("a:lstStyle"))
    etree.SubElement(lstStyle, qn("a:lvl1pPr"))
    etree.SubElement(txBody, qn("a:p"))


# ── Main builder ───────────────────────────────────────────────────────────────


def build_template(output_path: Path) -> dict:
    """Build the template PPTX and return the config dict."""
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    master = prs.slide_master
    all_layouts = list(master.slide_layouts)

    config_layouts: dict = {}

    for spec in LAYOUTS:
        layout = all_layouts[spec["layout_idx"]]

        # 1. Rename
        cSld = layout._element.find(qn("p:cSld"))
        cSld.set("name", spec["name"])

        # 2. Background colour
        bg_fill = layout.background.fill
        bg_fill.solid()
        bg_fill.fore_color.rgb = spec["bg"]

        spTree = cSld.find(qn("p:spTree"))

        # 3. Remove unwanted placeholders
        _remove_date_footer_slidenum(spTree)
        for extra_idx in spec["extra_phs"]:
            _remove_placeholder(spTree, extra_idx)

        # 4. If no body needed, remove it
        if spec["body"] is None:
            _remove_placeholder(spTree, 1)

        # 5. For layouts like "Blank" that have no placeholders, add them
        phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}

        if 0 not in phs:
            tx, ty, tw, th = spec["title"][:4]
            _add_placeholder(spTree, 0, tx, ty, tw, th)

        if spec["body"] is not None and 1 not in phs:
            bx, by, bw, bh = spec["body"][:4]
            _add_placeholder(spTree, 1, bx, by, bw, bh)

        # Re-fetch placeholders after additions
        phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}

        # 6. Reposition + restyle title placeholder
        if 0 in phs:
            tx, ty, tw, th, tpt, tbold, tcolor = spec["title"]
            _set_ph_position(phs[0]._element, tx, ty, tw, th)
            _set_ph_font(phs[0]._element, tpt, tbold, tcolor)

        # 7. Reposition + restyle body placeholder
        if spec["body"] is not None:
            phs = {ph.placeholder_format.idx: ph for ph in layout.placeholders}
            if 1 in phs:
                bx, by, bw, bh, bpt, bbold, bcolor = spec["body"]
                _set_ph_position(phs[1]._element, bx, by, bw, bh)
                _set_ph_font(phs[1]._element, bpt, bbold, bcolor)

        # 8. Add decorative shapes (insert BEFORE placeholders in z-order)
        for deco_i, (dx, dy, dw, dh, dc) in enumerate(spec["deco"]):
            elem_id = 500 + spec["layout_idx"] * 10 + deco_i
            _add_rect_shape(spTree, elem_id, dx, dy, dw, dh, dc)
            inserted = spTree[-1]
            spTree.remove(inserted)
            spTree.insert(2 + deco_i, inserted)

        # 9. Record config
        config_layouts[spec["name"]] = {
            "title_ph": 0,
            "body_ph": 1 if spec["body"] is not None else None,
        }

    prs.save(str(output_path))
    return config_layouts


def main() -> None:
    """Builds the default template for SlideForge."""
    out_dir = Path("slideforge/templates/default")
    out_dir.mkdir(parents=True, exist_ok=True)

    slides_path = out_dir / "slides.pptx"
    config_path = out_dir / "config.json"

    print("Building template...")
    config_layouts = build_template(slides_path)

    config = {"version": 1, "layouts": config_layouts}
    config_path.write_text(json.dumps(config, indent=2, ensure_ascii=False), encoding="utf-8")

    print(f"  {slides_path}  ({slides_path.stat().st_size // 1024} KB)")
    print(f"  {config_path}")
    print()
    print("Layout summary:")
    for name, cfg in config_layouts.items():
        body = f"PH {cfg['body_ph']}" if cfg["body_ph"] is not None else "—"
        print(f"  {name:20s}  title=PH {cfg['title_ph']}  body={body}")


if __name__ == "__main__":
    main()
