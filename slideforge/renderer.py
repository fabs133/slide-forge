"""PPTX renderer — converts a Presentation model into a .pptx file."""
from __future__ import annotations

import re
from pathlib import Path

from lxml import etree
from pptx import Presentation as PptxPresentation
from pptx.oxml.ns import qn

from .models import Presentation, Slide
from .template_loader import TemplateLoader, load_template

_DEFAULT_TEMPLATE = Path(__file__).parent / "templates" / "default"

_BULLET_PREFIX = re.compile(r"^[-•]\s+")

# ── Dynamic font sizing ──────────────────────────────────────────────────────

# (min_pt, max_pt) per layout — few bullets → big text, many → small text
BODY_FONT_LIMITS: dict[str, tuple[int, int]] = {
    "SP_Title": (16, 22),
    "SP_Content": (10, 28),
    "SP_Intro": (10, 28),
    "SP_Closing": (10, 28),
    "SP_Sources": (8, 16),
    "SP_SectionBreak": (16, 40),
    "SP_Code": (8, 16),
}

LINE_SPACING = 1.25


class RendererError(Exception):
    """Raised when rendering fails."""


def compute_font_size(
    num_lines: int,
    box_height_emu: int,
    min_pt: int,
    max_pt: int,
) -> float:
    """Compute an optimal font size (in pt) to fill the box without overflow."""
    if num_lines == 0:
        return max_pt
    box_height_pt = box_height_emu / 12700
    size = box_height_pt / (num_lines * LINE_SPACING)
    return max(min_pt, min(max_pt, size))


def _apply_dynamic_font_size(placeholder, layout_name: str, num_lines: int) -> None:
    """Set explicit font size on every run in the body placeholder."""
    limits = BODY_FONT_LIMITS.get(layout_name)
    if limits is None:
        return

    # Use python-pptx .height which resolves layout inheritance
    box_height = placeholder.height
    if not box_height:
        return

    min_pt, max_pt = limits
    size_pt = compute_font_size(num_lines, box_height, min_pt, max_pt)
    size_val = int(size_pt * 100)  # OOXML stores font size in hundredths of a point

    txBody = placeholder._element.find(qn("p:txBody"))
    if txBody is None:
        return

    # Remove any autofit elements — we control sizing explicitly
    bodyPr = txBody.find(qn("a:bodyPr"))
    if bodyPr is not None:
        for tag in ("a:normAutofit", "a:spAutoFit"):
            el = bodyPr.find(qn(tag))
            if el is not None:
                bodyPr.remove(el)

    for para in txBody.findall(qn("a:p")):
        for run in para.findall(qn("a:r")):
            rPr = run.find(qn("a:rPr"))
            if rPr is None:
                rPr = etree.SubElement(run, qn("a:rPr"))
                # rPr must be first child of r
                run.remove(rPr)
                run.insert(0, rPr)
            rPr.set("sz", str(size_val))


def render_pptx(
    presentation: Presentation,
    output: Path,
    template_dir: Path | None = None,
) -> None:
    """Render a Presentation model to a .pptx file.

    Args:
        presentation: The data model to render.
        output: Path to write the .pptx file.
        template_dir: Template directory (defaults to built-in default).
    """
    loader = load_template(template_dir or _DEFAULT_TEMPLATE)
    prs = loader.open_presentation()

    for slide_model in presentation.slides:
        _render_slide(prs, loader, slide_model)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))


def _render_slide(prs: PptxPresentation, loader: TemplateLoader, slide: Slide) -> None:
    """Render one Slide model onto the presentation."""
    layout_name = slide.layout

    if layout_name not in loader.layouts:
        from .template_loader import TemplateLayoutError

        raise TemplateLayoutError(f"Layout '{layout_name}' not in template config")

    tl = loader.layouts[layout_name]
    layout = loader.get_layout(prs, layout_name)
    pptx_slide = prs.slides.add_slide(layout)

    ph_map = {ph.placeholder_format.idx: ph for ph in pptx_slide.placeholders}

    # Title placeholder (always present)
    if tl.title_ph in ph_map:
        ph_map[tl.title_ph].text = slide.title

    # Body placeholder (None for SP_SectionBreak)
    if tl.body_ph is not None and tl.body_ph in ph_map:
        body_ph = ph_map[tl.body_ph]
        tf = body_ph.text_frame
        tf.clear()
        lines = [_BULLET_PREFIX.sub("", line) for line in slide.body.split("\n") if line.strip()]
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
        _apply_dynamic_font_size(body_ph, layout_name, len(lines))

    # Speaker notes
    if slide.notes.strip():
        pptx_slide.notes_slide.notes_text_frame.text = slide.notes
